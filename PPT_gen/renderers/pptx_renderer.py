import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

C_BLACK   = RGBColor(0x1A, 0x1A, 0x1A)  
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY_BG = RGBColor(0xEB, 0xEB, 0xEB)  
C_GREY_FT = RGBColor(0x55, 0x55, 0x55)  
C_NUM     = RGBColor(0x1A, 0x1A, 0x1A)  

FONT_SERIF = "Georgia"        
FONT_SANS  = "Calibri"        

W  = Inches(13.33)
H  = Inches(7.5)
ML = Inches(0.5)
MR = Inches(0.5)
MT = Inches(0.5)

TITLE_Y   = Inches(0.38)
TITLE_H   = Inches(1.1)           
CONTENT_Y = Inches(1.6)           
CONTENT_H = H - CONTENT_Y - Inches(0.7)


FOOTER_Y  = H - Inches(0.55)
FOOTER_H  = Inches(0.4)


PANEL_X   = Inches(6.7)
PANEL_W   = W - PANEL_X
PANEL_H   = H




def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _tb(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.fill.background()
    tb.line.fill.background()
    return tb


def _para(slide, x, y, w, h, text, size_pt, color,
          bold=False, italic=False,
          align=PP_ALIGN.LEFT, wrap=True,
          spacing_pt=None, font=FONT_SERIF):
    tb = _tb(slide, x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    if spacing_pt:
        p.line_spacing = Pt(spacing_pt)
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size_pt)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.italic = italic
    return tb


import re

def _bullet_para(tf, text, first=False, size_pt=14,
                 space_before_pt=10, color=C_BLACK):

    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment    = PP_ALIGN.LEFT
    p.space_before = Pt(space_before_pt)
    p.space_after  = Pt(0)
    p.line_spacing = Pt(size_pt + 6)

    pPr = p._p.get_or_add_pPr()
    pPr.set('indent', '-216000')
    pPr.set('marL',    '216000')
    etree.SubElement(pPr, qn('a:buNone'))

    dash_run = p.add_run()
    dash_run.text = "—  "
    dash_run.font.name  = FONT_SERIF
    dash_run.font.size  = Pt(size_pt)
    dash_run.font.color.rgb = color
    dash_run.font.bold  = False

    segments = re.split(r'\*\*(.*?)\*\*', text)
    for i, segment in enumerate(segments):
        if not segment:
            continue
        run = p.add_run()
        run.text = segment
        run.font.name  = FONT_SERIF
        run.font.size  = Pt(size_pt)
        run.font.color.rgb = color
        run.font.bold  = (i % 2 == 1) 
    return p




def _footer(slide, metadata):
    short  = metadata.project_short_name or metadata.project_name[:50]
    client = metadata.client_name

    cx = W / 2 - Inches(3)
    cw = Inches(6)

    tb = _tb(slide, cx, FOOTER_Y - Inches(0.02), cw, FOOTER_H)
    tf = tb.text_frame
    tf.word_wrap = False

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = short
    r1.font.name  = FONT_SERIF
    r1.font.size  = Pt(9)
    r1.font.bold  = True
    r1.font.color.rgb = C_GREY_FT

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = client
    r2.font.name  = FONT_SERIF
    r2.font.size  = Pt(9)
    r2.font.bold  = False
    r2.font.color.rgb = C_GREY_FT


def _page_number(slide, num: int):
    _para(slide,
          W - Inches(0.75), FOOTER_Y, Inches(0.55), Inches(0.3),
          str(num), 9, C_GREY_FT,
          bold=False, align=PP_ALIGN.RIGHT,
          wrap=False, font=FONT_SANS)


def render_cover_slide(prs, metadata):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C_WHITE)

    _para(slide,
          0, Inches(0.28), W, Inches(0.45),
          "Perkins&Will",
          13, C_BLACK, bold=True,
          align=PP_ALIGN.CENTER, wrap=False)

    name = metadata.project_name
    fs   = 52 if len(name) <= 30 else 40 if len(name) <= 55 else 30
    _para(slide,
          Inches(1.2), Inches(2.6),
          W - Inches(2.4), Inches(2.4),
          name,
          fs, C_BLACK, bold=True,
          align=PP_ALIGN.CENTER, wrap=True,
          spacing_pt=fs + 8)

    _para(slide,
          Inches(1.2), Inches(6.15),
          W - Inches(2.4), Inches(0.45),
          f"Prepared for: {metadata.client_name}",
          13, C_BLACK, bold=True,
          align=PP_ALIGN.CENTER, wrap=False)

    _para(slide,
          Inches(1.2), Inches(6.6),
          W - Inches(2.4), Inches(0.35),
          "RFP Analysis",
          11, C_GREY_FT, bold=False, italic=True,
          align=PP_ALIGN.CENTER, wrap=False,
          font=FONT_SANS)

    if metadata.client_logo_path and os.path.exists(metadata.client_logo_path):
        try:
            slide.shapes.add_picture(
                metadata.client_logo_path,
                W - Inches(2.4), Inches(0.18),
                width=Inches(2.0)
            )
        except Exception:
            pass



def _layout_full(slide, sd, metadata, idx):
    _rect(slide, 0, 0, W, H, C_WHITE)

    title = sd.title or ""
    title_fs = 24 if len(title) > 35 else 28
    _para(slide,
          ML, TITLE_Y, W - ML - MR, TITLE_H,
          title, title_fs, C_BLACK, bold=True,
          wrap=True, spacing_pt=title_fs + 4)

    bullets = sd.bullets or []
    if not bullets:
        _footer(slide, metadata)
        _page_number(slide, idx)
        return

    n    = len(bullets)
    size = 15 if n <= 4 else 14 if n <= 7 else 13
    ch_pt        = CONTENT_H / 12700 * 72
    line_h_pt    = size + 6
    total_h_pt   = n * (line_h_pt + 10)
    slack        = max(0, ch_pt - total_h_pt)
    space_before = min(int(slack / max(n, 1)), 20)

    if n >= 9:
        mid   = (n + 1) // 2
        col_w = (W - ML - MR) / 2 - Inches(0.2)
        for col, col_b in enumerate([bullets[:mid], bullets[mid:]]):
            cx = ML + col * (col_w + Inches(0.4))
            tb = _tb(slide, cx, CONTENT_Y, col_w, CONTENT_H)
            tf = tb.text_frame
            tf.word_wrap = True
            for i, b in enumerate(col_b):
                _bullet_para(tf, b, first=(i == 0),
                             size_pt=size, space_before_pt=space_before)
    else:
        tb = _tb(slide, ML, CONTENT_Y, W - ML - MR, CONTENT_H)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            _bullet_para(tf, b, first=(i == 0),
                         size_pt=size, space_before_pt=space_before)

    _footer(slide, metadata)
    _page_number(slide, idx)



def _layout_with_image(slide, sd, metadata, idx):
    img_path = sd.image_path or sd.ai_image_path
    if not (img_path and os.path.exists(img_path)):
        _layout_full(slide, sd, metadata, idx)
        return

    _rect(slide, 0, 0, W, H, C_WHITE)

    _rect(slide, PANEL_X, 0, PANEL_W, PANEL_H, C_GREY_BG)

    title_txt_w = PANEL_X - ML - Inches(0.3)
    title_size  = 24 if len(sd.title or "") > 30 else 26
    _para(slide,
          ML, TITLE_Y, title_txt_w, TITLE_H,
          sd.title or "", title_size, C_BLACK, bold=True,
          wrap=True, spacing_pt=title_size + 4)

    try:
        slide.shapes.add_picture(
            img_path, PANEL_X, 0, PANEL_W, PANEL_H)
    except Exception:
        _layout_full(slide, sd, metadata, idx)
        return

    bullets = sd.bullets or []
    if bullets:
        n    = len(bullets)
        size = 14 if n <= 5 else 13
        txt_w = PANEL_X - ML - Inches(0.3)
        tb = _tb(slide, ML, CONTENT_Y, txt_w, CONTENT_H)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            _bullet_para(tf, b, first=(i == 0),
                         size_pt=size, space_before_pt=9)

    _footer(slide, metadata)
    _page_number(slide, idx)



def _layout_image_dominant(slide, sd, metadata, idx):
    img_path = sd.image_path or sd.ai_image_path
    has_img  = bool(img_path and os.path.exists(img_path))

    _rect(slide, 0, 0, W, H, C_WHITE)

    if has_img:
        img_x = Inches(5.5)
        img_w = W - img_x
        try:
            _rect(slide, img_x, 0, img_w, H, C_GREY_BG)
            slide.shapes.add_picture(img_path, img_x, 0, img_w, H)
        except Exception:
            has_img = False

    txt_w = Inches(5.0) if has_img else W - ML - MR

    _para(slide,
          ML, TITLE_Y, txt_w, TITLE_H,
          sd.title or "", 26, C_BLACK, bold=True,
          wrap=True, spacing_pt=30)

    bullets = (sd.bullets or [])[:5]
    if bullets:
        n    = len(bullets)
        size = 14
        tb   = _tb(slide, ML, CONTENT_Y, txt_w, CONTENT_H)
        tf   = tb.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            _bullet_para(tf, b, first=(i == 0),
                         size_pt=size, space_before_pt=10)

    _footer(slide, metadata)
    _page_number(slide, idx)



def _layout_agenda(slide, sd, metadata, idx):
    _rect(slide, 0, 0, W, H, C_WHITE)
    _rect(slide, PANEL_X, 0, PANEL_W, PANEL_H, C_GREY_BG)

    img_path = sd.image_path or sd.ai_image_path
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, PANEL_X, 0, PANEL_W, PANEL_H)
        except Exception:
            pass

    txt_w = PANEL_X - ML - Inches(0.3)

    _para(slide,
          ML, TITLE_Y, txt_w, TITLE_H,
          sd.title or "", 32, C_BLACK, bold=True,
          wrap=False)

    bullets = sd.bullets or []
    by = CONTENT_Y + Inches(0.15)
    for i, b in enumerate(bullets[:8]):
        num_str = f"{i+1:02d} \u2014 "
        line_h  = Inches(0.55)

        tb = _tb(slide, ML, by, txt_w, line_h)
        tf = tb.text_frame
        tf.word_wrap = False
        p  = tf.paragraphs[0]
        p.space_before = Pt(0)
        p.line_spacing = Pt(20)

        r1 = p.add_run()
        r1.text = num_str
        r1.font.name  = FONT_SERIF
        r1.font.size  = Pt(16)
        r1.font.bold  = True
        r1.font.color.rgb = C_BLACK

        r2 = p.add_run()
        r2.text = b
        r2.font.name  = FONT_SERIF
        r2.font.size  = Pt(16)
        r2.font.bold  = True
        r2.font.color.rgb = C_BLACK

        by += line_h + Inches(0.05)

    _footer(slide, metadata)
    _page_number(slide, idx)



def _pick_layout(sd, idx: int) -> str:
    img_path = sd.image_path or sd.ai_image_path
    has_img  = bool(img_path and os.path.exists(img_path))
    desc     = (sd.image_description or "").lower()
    title    = (sd.title or "").lower()

    if has_img and any(k in desc for k in [
        "map", "aerial", "campus", "boundary", "plan", "site",
        "floor plan", "elevation", "district"
    ]):
        return "image_dominant"

    if any(k in title for k in ["agenda", "overview", "contents", "table"]):
        return "agenda"

    if has_img:
        return "with_image"

    return "full"



def _resolve_image_path(path: str) -> str:
    if not path:
        return path
    if os.path.exists(path):
        return path

    basename = os.path.basename(path)
    search_roots = [
        Path.cwd(),
        Path.cwd() / "images",
        Path.cwd() / "tmp",
        Path.cwd() / "tmp" / "debug" / "images",
        Path(__file__).resolve().parent.parent.parent / "Docint" / "tmp" / "debug" / "images",
    ]
    for root in search_roots:
        for candidate in root.glob("**/" + basename):
            return str(candidate)
    return path



def render_pptx(deck, output_path):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    if deck.metadata.client_logo_path:
        deck.metadata.client_logo_path = _resolve_image_path(
            deck.metadata.client_logo_path)
    for sd in deck.slides:
        if sd.image_path:
            sd.image_path = _resolve_image_path(sd.image_path)
        if sd.ai_image_path:
            sd.ai_image_path = _resolve_image_path(sd.ai_image_path)

    render_cover_slide(prs, deck.metadata)

    for i, sd in enumerate(deck.slides):
        if "cover" in (sd.title or "").lower():
            continue

        slide  = prs.slides.add_slide(prs.slide_layouts[6])
        layout = _pick_layout(sd, i)

        if layout == "image_dominant":
            _layout_image_dominant(slide, sd, deck.metadata, i + 1)
        elif layout == "agenda":
            _layout_agenda(slide, sd, deck.metadata, i + 1)
        elif layout == "with_image":
            _layout_with_image(slide, sd, deck.metadata, i + 1)
        else:
            _layout_full(slide, sd, deck.metadata, i + 1)

    prs.save(output_path)
    print(f"  [Renderer] Saved → {output_path}")